import pandas as pd
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from io import BytesIO

# Cargar datos
def load_data():
    df = pd.read_csv("supermarket_Sales.csv")
    columnas_traducidas = {
        'Invoice ID': 'ID de Factura',
        'Branch': 'Sucursal',
        'City': 'Ciudad',
        'Customer type': 'Tipo de Cliente',
        'Gender': 'Género',
        'Product line': 'Línea de Producto',
        'Unit price': 'Precio Unitario',
        'Quantity': 'Cantidad',
        'Tax 5%': 'Impuesto 5%',
        'Total': 'Total',
        'Date': 'Fecha',
        'Time': 'Hora',
        'Payment': 'Método de Pago',
        'Cost of goods sold': 'Costo de Bienes Vendidos',
        'Gross margin percentage': 'Porcentaje de Margen Bruto',
        'Gross income': 'Ingreso Bruto',
        'Customer stratification rating': 'Calificación de Estratificación del Cliente'
    }
    df = df.rename(columns=columnas_traducidas)
    df['Fecha'] = pd.to_datetime(df['Fecha'])
    return df

df = load_data()

# Crear presentación
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Función mejorada para agregar una diapositiva con un gráfico
def add_chart_slide(prs, title, fig, layout_index=5):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    slide.shapes.title.text = title
    
    # Ajustar el tamaño de la figura
    fig.update_layout(width=1000, height=600)
    
    img_bytes = fig.to_image(format="png")
    img_stream = BytesIO(img_bytes)
    
    # Centrar la imagen en la diapositiva
    left = (prs.slide_width - Inches(10)) / 2
    top = (prs.slide_height - Inches(6)) / 2
    slide.shapes.add_picture(img_stream, left, top, width=Inches(10), height=Inches(6))

# Diapositiva de título
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Dashboard de Ventas de Supermercado"

# Diapositiva de métricas mejorada
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Métricas Principales"
metrics = [
    ("Ventas Totales", f"${df['Total'].sum():,.2f}"),
    ("Transacciones", f"{len(df):,}"),
    ("Venta Promedio", f"${df['Total'].mean():,.2f}"),
    ("Margen Bruto Promedio", f"${df['Ingreso Bruto'].mean():,.2f}")
]

for i, (name, value) in enumerate(metrics):
    left = Inches(1 + i * 3.75)
    top = Inches(2)
    width = Inches(3.5)
    height = Inches(2)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(225, 225, 225)  # Light gray background
    shape.line.color.rgb = RGBColor(200, 200, 200)  # Lighter border
    
    tf = shape.text_frame
    tf.text = f"{name}\n{value}"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[1].font.size = Pt(24)
    tf.paragraphs[1].font.color.rgb = RGBColor(0, 112, 192)  # Blue color for value

# Ventas Diarias
ventas_diarias = df.groupby('Fecha')['Total'].sum().reset_index()
fig_ventas_diarias = px.line(ventas_diarias, x='Fecha', y='Total', title='Ventas Diarias')
add_chart_slide(prs, "Ventas Diarias", fig_ventas_diarias)

# Ventas por Tipo de Cliente y Género
ventas_tipo_genero = df.groupby(['Tipo de Cliente', 'Género'])['Total'].sum().unstack()
fig_tipo_genero = px.bar(ventas_tipo_genero, title='Ventas por Tipo de Cliente y Género', barmode='group')
add_chart_slide(prs, "Ventas por Tipo de Cliente y Género", fig_tipo_genero)

# Ventas por Línea de Producto
productos = df['Línea de Producto'].value_counts()
fig_productos = px.pie(values=productos.values, names=productos.index, title='Ventas por Línea de Producto')
add_chart_slide(prs, "Ventas por Línea de Producto", fig_productos)

# Ventas por Ciudad
ventas_ciudad = df.groupby('Ciudad')['Total'].sum().sort_values(ascending=True)
fig_ciudad = px.bar(ventas_ciudad, orientation='h', title='Ventas por Ciudad')
add_chart_slide(prs, "Ventas por Ciudad", fig_ciudad)

# Métodos de Pago
metodos_pago = df['Método de Pago'].value_counts()
fig_metodos_pago = px.pie(values=metodos_pago.values, names=metodos_pago.index, title='Métodos de Pago')
add_chart_slide(prs, "Métodos de Pago", fig_metodos_pago)

# Guardar presentación
prs.save('dashboard_ventas_supermercado.pptx')

print("Presentación de PowerPoint creada: dashboard_ventas_supermercado.pptx")